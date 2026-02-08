import io
import json
import os
import re
from typing import Any

from reportlab.lib.pagesizes import LETTER
from reportlab.pdfbase.pdfmetrics import stringWidth
from reportlab.pdfgen import canvas
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from app.course_jobs import get_job


def load_course(job_id: str) -> dict | None:
    job = get_job(job_id)
    if job and job.get("result"):
        return job["result"]

    file_path = os.path.join("data", "courses", f"{job_id}.json")
    if os.path.exists(file_path):
        with open(file_path, "r", encoding="utf-8") as f:
            return json.load(f)
    return None


def _slug(value: str) -> str:
    cleaned = re.sub(r"[^A-Za-z0-9]+", "-", value).strip("-")
    if not cleaned:
        return "course"
    return cleaned[:60]


def _wrap_text(text: str, max_width: float, font_name: str, font_size: int) -> list[str]:
    words = text.split()
    if not words:
        return [""]
    lines: list[str] = []
    current = words[0]
    for word in words[1:]:
        test_line = f"{current} {word}"
        if stringWidth(test_line, font_name, font_size) <= max_width:
            current = test_line
        else:
            lines.append(current)
            current = word
    lines.append(current)
    return lines


def _draw_paragraph(c: canvas.Canvas, text: str, x: float, y: float, width: float,
                    font_name: str = "Times-Roman", font_size: int = 11,
                    leading: int = 15, bottom_margin: int = 60) -> float:
    lines = _wrap_text(text, width, font_name, font_size)
    c.setFont(font_name, font_size)
    for line in lines:
        if y < bottom_margin:
            c.showPage()
            y = LETTER[1] - bottom_margin
            c.setFont(font_name, font_size)
        c.drawString(x, y, line)
        y -= leading
    return y


def _draw_heading(c: canvas.Canvas, text: str, x: float, y: float, width: float,
                  level: int = 1, bottom_margin: int = 60) -> float:
    sizes = {1: 18, 2: 14, 3: 12}
    font_size = sizes.get(level, 12)
    if y < bottom_margin + 20:
        c.showPage()
        y = LETTER[1] - bottom_margin
    c.setFont("Helvetica-Bold", font_size)
    return _draw_paragraph(c, text, x, y, width, font_name="Helvetica-Bold",
                           font_size=font_size, leading=font_size + 6,
                           bottom_margin=bottom_margin) - 6


def _draw_markdown(c: canvas.Canvas, text: str, x: float, y: float, width: float,
                   bottom_margin: int = 60) -> float:
    normalized = _normalize_markdown(text)
    lines = [line.rstrip() for line in normalized.splitlines()]
    for raw in lines:
        line = raw.strip()
        if not line:
            y -= 8
            continue
        if line.startswith("### "):
            y = _draw_heading(c, line[4:], x, y, width, level=3, bottom_margin=bottom_margin)
            continue
        if line.startswith("## "):
            y = _draw_heading(c, line[3:], x, y, width, level=2, bottom_margin=bottom_margin)
            continue
        if line.startswith("# "):
            y = _draw_heading(c, line[2:], x, y, width, level=1, bottom_margin=bottom_margin)
            continue
        if line.startswith("- "):
            y = _draw_paragraph(c, f"• {line[2:]}", x + 10, y, width - 10,
                                font_name="Times-Roman", font_size=11, leading=15,
                                bottom_margin=bottom_margin)
            continue
        if re.match(r"^\d+\.\s", line):
            y = _draw_paragraph(c, line, x + 6, y, width - 6,
                                font_name="Times-Roman", font_size=11, leading=15,
                                bottom_margin=bottom_margin)
            continue
        y = _draw_paragraph(c, line, x, y, width, font_name="Times-Roman",
                            font_size=11, leading=15, bottom_margin=bottom_margin)
    return y


def _normalize_markdown(text: str) -> str:
    if not text:
        return ""
    # Ensure headings start on new lines
    text = re.sub(r"\s*(#{1,3}\s)", r"\n\1", text)
    # Ensure list items start on new lines
    text = re.sub(r"\s*(-\s+)", r"\n\1", text)
    text = re.sub(r"\s*(\d+\.\s+)", r"\n\1", text)
    # Collapse excessive blank lines
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def build_course_pdf(course: dict) -> bytes:
    buffer = io.BytesIO()
    c = canvas.Canvas(buffer, pagesize=LETTER)
    width, height = LETTER
    margin_x = 54
    y = height - 72

    title = course.get("course_title") or "Course"
    hook = course.get("hook") or ""
    difficulty = course.get("difficulty") or "unknown"
    total_minutes = course.get("estimated_total_minutes") or 0
    source = course.get("source") or {}

    c.setFont("Helvetica-Bold", 22)
    c.drawString(margin_x, y, title)
    y -= 28

    if hook:
        y = _draw_paragraph(c, hook, margin_x, y, width - 2 * margin_x, font_name="Times-Italic",
                            font_size=12, leading=16)
        y -= 8

    c.setFont("Helvetica", 12)
    c.drawString(margin_x, y, f"Difficulty: {difficulty}")
    y -= 16
    c.drawString(margin_x, y, f"Estimated minutes: {total_minutes}")
    y -= 16
    if source.get("playlist_url"):
        y = _draw_paragraph(c, f"Source playlist: {source.get('playlist_url')}",
                            margin_x, y, width - 2 * margin_x, font_size=10, leading=12)
        y -= 8

    modules = course.get("modules") or []
    for module_index, module in enumerate(modules, start=1):
        if y < 120:
            c.showPage()
            y = height - 72

        c.setFont("Helvetica-Bold", 16)
        c.drawString(margin_x, y, f"Module {module_index}: {module.get('title')}")
        y -= 22

        objectives = module.get("objectives") or []
        if objectives:
            c.setFont("Helvetica-Bold", 12)
            c.drawString(margin_x, y, "Objectives")
            y -= 16
            for obj in objectives:
                y = _draw_paragraph(c, f"- {obj}", margin_x + 8, y,
                                    width - 2 * margin_x - 8, font_size=11)
            y -= 6

        module_minutes = module.get("estimated_minutes") or 0
        c.setFont("Helvetica", 11)
        c.drawString(margin_x, y, f"Estimated minutes: {module_minutes}")
        y -= 18

        lessons = module.get("lessons") or []
        for lesson_index, lesson in enumerate(lessons, start=1):
            if y < 120:
                c.showPage()
                y = height - 72

            c.setFont("Helvetica-Bold", 13)
            c.drawString(margin_x, y, f"Lesson {module_index}.{lesson_index}: {lesson.get('title')}")
            y -= 18

            lesson_summary = lesson.get("summary") or ""
            if lesson_summary:
                y = _draw_paragraph(c, lesson_summary, margin_x, y, width - 2 * margin_x,
                                    font_name="Times-Roman", font_size=11, leading=15)
                y -= 8

            if lesson.get("video_url"):
                y = _draw_paragraph(c, f"Video: {lesson.get('video_url')}",
                                    margin_x, y, width - 2 * margin_x, font_size=10, leading=12)
                y -= 6

            lesson_meta = (
                f"Difficulty: {lesson.get('difficulty') or 'unknown'} | "
                f"Estimated minutes: {lesson.get('estimated_minutes') or 0}"
            )
            y = _draw_paragraph(c, lesson_meta, margin_x, y, width - 2 * margin_x, font_size=10, leading=12)
            y -= 6

            learning_objectives = lesson.get("learning_objectives") or []
            if learning_objectives:
                c.setFont("Helvetica-Bold", 11)
                c.drawString(margin_x, y, "Learning objectives")
                y -= 14
                for obj in learning_objectives:
                    y = _draw_paragraph(c, f"- {obj}", margin_x + 8, y,
                                        width - 2 * margin_x - 8, font_size=10, leading=12)
                y -= 4

            study_material = lesson.get("study_material_markdown") or ""
            if study_material:
                y = _draw_heading(c, "Study material", margin_x, y, width - 2 * margin_x, level=2)
                y = _draw_markdown(c, study_material, margin_x, y, width - 2 * margin_x)
                y -= 8

            reading_guide = lesson.get("reading_guide_markdown") or ""
            if reading_guide:
                y = _draw_heading(c, "Reading guide", margin_x, y, width - 2 * margin_x, level=3)
                y = _draw_markdown(c, reading_guide, margin_x, y, width - 2 * margin_x)
                y -= 4

        quiz = module.get("quiz") or []
        if quiz:
            if y < 120:
                c.showPage()
                y = height - 72
            c.setFont("Helvetica-Bold", 13)
            c.drawString(margin_x, y, "Quiz")
            y -= 18
            for question_index, item in enumerate(quiz, start=1):
                question = item.get("question") or ""
                y = _draw_paragraph(c, f"{question_index}. {question}",
                                    margin_x, y, width - 2 * margin_x)
                options = item.get("options") or []
                for option_index, option in enumerate(options, start=1):
                    y = _draw_paragraph(c, f"   {option_index}) {option}",
                                        margin_x, y, width - 2 * margin_x, font_size=10, leading=12)
                answer_index = item.get("answer_index")
                explanation = item.get("explanation") or ""
                if answer_index is not None:
                    y = _draw_paragraph(c, f"Answer: {answer_index + 1}",
                                        margin_x, y, width - 2 * margin_x, font_size=10, leading=12)
                if explanation:
                    y = _draw_paragraph(c, f"Explanation: {explanation}",
                                        margin_x, y, width - 2 * margin_x, font_size=10, leading=12)
                y -= 6

    c.save()
    buffer.seek(0)
    return buffer.read()


def _add_bullets(frame, items: list[str], max_chars: int = 200):
    for item in items:
        if not item:
            continue
        text = item.strip()
        if len(text) > max_chars:
            text = f"{text[:max_chars].rstrip()}..."
        p = frame.add_paragraph()
        p.text = text
        p.level = 0


def _apply_slide_theme(slide, prs, *, title_color: RGBColor, body_color: RGBColor,
                       background: RGBColor, accent: RGBColor):
    width = prs.slide_width
    height = prs.slide_height
    bg = slide.shapes.add_shape(
        1, 0, 0, width, height  # MSO_SHAPE.RECTANGLE = 1
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = background
    bg.line.fill.background()

    bar = slide.shapes.add_shape(1, 0, height - Inches(0.3), width, Inches(0.3))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    # Send background to back, bar above it
    sp_tree = slide.shapes._spTree
    sp_tree.remove(bg._element)
    sp_tree.insert(0, bg._element)
    sp_tree.remove(bar._element)
    sp_tree.insert(1, bar._element)

    if slide.shapes.title:
        title = slide.shapes.title
        for paragraph in title.text_frame.paragraphs:
            if paragraph.runs:
                for run in paragraph.runs:
                    run.font.color.rgb = title_color
                    run.font.bold = True
                    run.font.size = Pt(32)
            else:
                paragraph.font.color.rgb = title_color
                paragraph.font.bold = True
                paragraph.font.size = Pt(32)
        title.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

    for shape in slide.shapes:
        if not hasattr(shape, "text_frame") or shape == slide.shapes.title:
            continue
        for paragraph in shape.text_frame.paragraphs:
            if paragraph.runs:
                for run in paragraph.runs:
                    run.font.color.rgb = body_color
                    run.font.size = Pt(18)
            else:
                paragraph.font.color.rgb = body_color
                paragraph.font.size = Pt(18)


def build_course_pptx(course: dict) -> bytes:
    prs = Presentation()
    theme = {
        "background": RGBColor(15, 23, 42),
        "title": RGBColor(248, 250, 252),
        "body": RGBColor(226, 232, 240),
        "accent": RGBColor(99, 102, 241)
    }

    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = course.get("course_title") or "Course"
    subtitle = title_slide.placeholders[1]
    subtitle.text = course.get("hook") or "Course overview"
    _apply_slide_theme(
        title_slide,
        prs,
        title_color=theme["title"],
        body_color=theme["body"],
        background=theme["background"],
        accent=theme["accent"]
    )

    overview = prs.slides.add_slide(prs.slide_layouts[1])
    overview.shapes.title.text = "Course Overview"
    overview_frame = overview.shapes.placeholders[1].text_frame
    overview_frame.clear()
    overview_frame.text = f"Difficulty: {course.get('difficulty') or 'unknown'}"
    p = overview_frame.add_paragraph()
    p.text = f"Estimated minutes: {course.get('estimated_total_minutes') or 0}"
    source = course.get("source") or {}
    if source.get("playlist_url"):
        p = overview_frame.add_paragraph()
        p.text = f"Playlist: {source.get('playlist_url')}"
    _apply_slide_theme(
        overview,
        prs,
        title_color=theme["title"],
        body_color=theme["body"],
        background=theme["background"],
        accent=theme["accent"]
    )

    modules = course.get("modules") or []
    for module_index, module in enumerate(modules, start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = f"Module {module_index}: {module.get('title')}"
        frame = slide.shapes.placeholders[1].text_frame
        frame.clear()
        frame.text = f"Estimated minutes: {module.get('estimated_minutes') or 0}"
        _add_bullets(frame, module.get("objectives") or [])
        _apply_slide_theme(
            slide,
            prs,
            title_color=theme["title"],
            body_color=theme["body"],
            background=theme["background"],
            accent=theme["accent"]
        )

        lessons = module.get("lessons") or []
        for lesson_index, lesson in enumerate(lessons, start=1):
            lesson_slide = prs.slides.add_slide(prs.slide_layouts[1])
            lesson_slide.shapes.title.text = f"Lesson {module_index}.{lesson_index}: {lesson.get('title')}"
            lesson_frame = lesson_slide.shapes.placeholders[1].text_frame
            lesson_frame.clear()
            lesson_frame.text = lesson.get("summary") or "Lesson summary"

            objectives = lesson.get("learning_objectives") or []
            if objectives:
                _add_bullets(lesson_frame, [f"Objective: {obj}" for obj in objectives], max_chars=180)

            meta_parts = [
                f"Difficulty: {lesson.get('difficulty') or 'unknown'}",
                f"Estimated minutes: {lesson.get('estimated_minutes') or 0}"
            ]
            _add_bullets(lesson_frame, [" | ".join(meta_parts)], max_chars=180)

            video_url = lesson.get("video_url")
            if video_url:
                _add_bullets(lesson_frame, [f"Video: {video_url}"], max_chars=180)

            study_material = lesson.get("study_material_markdown") or ""
            reading_guide = lesson.get("reading_guide_markdown") or ""
            notes_text = study_material or reading_guide
            if notes_text:
                notes = lesson_slide.notes_slide.notes_text_frame
                notes.text = notes_text[:2000]
            _apply_slide_theme(
                lesson_slide,
                prs,
                title_color=theme["title"],
                body_color=theme["body"],
                background=theme["background"],
                accent=theme["accent"]
            )

        quiz = module.get("quiz") or []
        if quiz:
            quiz_slide = prs.slides.add_slide(prs.slide_layouts[1])
            quiz_slide.shapes.title.text = f"Module {module_index} Quiz"
            quiz_frame = quiz_slide.shapes.placeholders[1].text_frame
            quiz_frame.clear()
            for question_index, item in enumerate(quiz, start=1):
                question = item.get("question") or ""
                if question:
                    p = quiz_frame.add_paragraph()
                    p.text = f"{question_index}. {question}"
                options = item.get("options") or []
                for option_index, option in enumerate(options, start=1):
                    p = quiz_frame.add_paragraph()
                    p.text = f"   {option_index}) {option}"
                answer_index = item.get("answer_index")
                if answer_index is not None:
                    p = quiz_frame.add_paragraph()
                    p.text = f"Answer: {answer_index + 1}"
                explanation = item.get("explanation") or ""
                if explanation:
                    p = quiz_frame.add_paragraph()
                    p.text = f"Explanation: {explanation}"
            _apply_slide_theme(
                quiz_slide,
                prs,
                title_color=theme["title"],
                body_color=theme["body"],
                background=theme["background"],
                accent=theme["accent"]
            )

    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output.read()


def build_export_filenames(course: dict) -> dict[str, str]:
    base = _slug(course.get("course_title") or "course")
    return {
        "pdf": f"{base}.pdf",
        "pptx": f"{base}.pptx"
    }
